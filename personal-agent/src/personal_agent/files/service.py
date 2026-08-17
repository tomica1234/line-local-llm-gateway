from __future__ import annotations

import mimetypes
import shutil
import zipfile
from pathlib import Path
from typing import Any
from uuid import uuid4

_SECRET_NAMES = {
    ".env",
    ".netrc",
    "credentials.json",
    "id_rsa",
    "id_ed25519",
    "known_hosts",
}
_SECRET_SUFFIXES = {".key", ".pem", ".p12", ".pfx", ".kdbx"}


class FileService:
    def __init__(self, roots: tuple[Path, ...], trash_root: Path):
        self.roots = tuple(root.resolve() for root in roots)
        self.trash_root = trash_root.resolve()

    def search(self, query: str, *, limit: int = 100) -> list[dict[str, object]]:
        normalized = query.casefold().strip()
        if not normalized or not self.roots:
            return []
        results: list[dict[str, object]] = []
        for root in self.roots:
            if not root.exists():
                continue
            for path in root.rglob("*"):
                if len(results) >= limit:
                    return results
                if not path.is_file() or self._is_secret(path):
                    continue
                if normalized in path.name.casefold():
                    results.append(self.metadata(path))
        return results

    def recent(self, *, limit: int = 100) -> list[dict[str, object]]:
        if not self.roots:
            return []
        candidates: list[Path] = []
        for root in self.roots:
            if not root.exists():
                continue
            for path in root.rglob("*"):
                if path.is_file() and not self._is_secret(path):
                    candidates.append(path)
        candidates.sort(key=lambda item: item.stat().st_mtime, reverse=True)
        return [self.metadata(path) for path in candidates[:limit]]

    def read(self, value: str, *, max_bytes: int = 1_000_000) -> dict[str, object]:
        path = self._resolve(value, must_exist=True)
        self._deny_secret(path)
        if not path.is_file():
            raise ValueError("Path is not a regular file")
        if path.stat().st_size > max_bytes:
            raise ValueError("File exceeds the bounded read size")
        media_type = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
        if not (
            media_type.startswith("text/")
            or path.suffix.lower() in {".json", ".md", ".csv", ".yaml", ".yml"}
        ):
            raise ValueError("Binary file extraction is not enabled for this type")
        content = path.read_text(encoding="utf-8", errors="replace")
        return {**self.metadata(path), "content": content, "trust_level": "untrusted"}

    def inspect(self, value: str) -> dict[str, object]:
        path = self._safe_file(value)
        suffix = path.suffix.casefold()
        archive = zipfile.is_zipfile(path)
        macro_capable = suffix in {".docm", ".xlsm", ".pptm"}
        supported = {
            ".pdf": ["extract_text", "extract_metadata"],
            ".docx": ["extract_text", "extract_metadata"],
            ".xlsx": ["extract_text", "extract_metadata"],
            ".pptx": ["extract_text", "extract_metadata"],
            ".html": ["extract_text", "extract_metadata"],
            ".htm": ["extract_text", "extract_metadata"],
            ".zip": ["list_archive", "extract_metadata"],
        }.get(suffix, ["extract_metadata"] if self._is_image(path) else [])
        return {
            **self.metadata(path),
            "format": suffix.removeprefix(".") or "unknown",
            "is_archive": archive,
            "macro_capable": macro_capable,
            "macros_executed": False,
            "supported_operations": supported,
        }

    def extract_metadata(self, value: str) -> dict[str, object]:
        path = self._safe_file(value)
        result: dict[str, object] = self.inspect(str(path))
        suffix = path.suffix.casefold()
        if suffix == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(path)
            result.update(
                {
                    "pages": len(reader.pages),
                    "encrypted": bool(reader.is_encrypted),
                    "document_metadata": {
                        str(key): str(value) for key, value in (reader.metadata or {}).items()
                    },
                }
            )
        elif self._is_image(path):
            from PIL import Image

            with Image.open(path) as image:
                result.update(
                    {
                        "width": image.width,
                        "height": image.height,
                        "image_format": image.format,
                        "mode": image.mode,
                        "frames": getattr(image, "n_frames", 1),
                    }
                )
        elif zipfile.is_zipfile(path):
            listing = self.list_archive(str(path), limit=1_000)
            result.update(
                {
                    "entries": listing["entry_count"],
                    "uncompressed_size": listing["uncompressed_size"],
                    "archive_safe": listing["safe"],
                }
            )
        return result

    def extract_text(
        self,
        value: str,
        *,
        pages: list[int] | None = None,
        max_chars: int = 500_000,
    ) -> dict[str, object]:
        path = self._safe_file(value)
        suffix = path.suffix.casefold()
        page_results: list[dict[str, object]] = []
        text = ""
        if suffix == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(path)
            if reader.is_encrypted:
                raise PermissionError("Encrypted PDF requires explicit human handling")
            indexes = self._page_indexes(len(reader.pages), pages)
            for index in indexes:
                page_text = reader.pages[index].extract_text() or ""
                page_results.append({"page": index + 1, "text": page_text, "tables": []})
            try:
                import fitz

                document = fitz.open(path)
                try:
                    for result, index in zip(page_results, indexes, strict=True):
                        finder = getattr(document[index], "find_tables", None)
                        if finder is not None:
                            result["tables"] = [table.extract() for table in finder().tables]
                finally:
                    document.close()
            except (RuntimeError, ValueError):
                # Native text remains available even when heuristic table detection fails.
                pass
            text = "\n\n".join(str(item["text"]) for item in page_results)
        elif suffix == ".docx":
            from docx import Document

            document = Document(path)
            text = "\n".join(paragraph.text for paragraph in document.paragraphs)
            for table in document.tables:
                for row in table.rows:
                    text += "\n" + "\t".join(cell.text for cell in row.cells)
        elif suffix == ".xlsx":
            from openpyxl import load_workbook

            workbook = load_workbook(path, read_only=True, data_only=True, keep_links=False)
            chunks: list[str] = []
            try:
                for sheet in workbook.worksheets:
                    chunks.append(f"# {sheet.title}")
                    for row in sheet.iter_rows(values_only=True):
                        chunks.append("\t".join("" if cell is None else str(cell) for cell in row))
            finally:
                workbook.close()
            text = "\n".join(chunks)
        elif suffix == ".pptx":
            from pptx import Presentation

            presentation = Presentation(path)
            chunks = []
            for number, slide in enumerate(presentation.slides, start=1):
                chunks.append(f"# Slide {number}")
                chunks.extend(
                    str(shape.text)
                    for shape in slide.shapes
                    if hasattr(shape, "text") and str(shape.text).strip()
                )
            text = "\n".join(chunks)
        elif suffix in {".html", ".htm"}:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="replace"), "html.parser")
            for node in soup(["script", "style", "template"]):
                node.decompose()
            text = soup.get_text("\n", strip=True)
        elif self._is_image(path):
            return {
                **self.extract_metadata(str(path)),
                "text": "",
                "vision_required": True,
                "fallback_reason": "IMAGE_HAS_NO_NATIVE_TEXT",
                "trust_boundary": "untrusted_external_content",
            }
        else:
            return self.read(str(path), max_bytes=min(max_chars * 4, 4_000_000))
        truncated = len(text) > max_chars
        text = text[:max_chars]
        return {
            **self.metadata(path),
            "text": text,
            "pages": page_results,
            "truncated": truncated,
            "vision_required": suffix == ".pdf" and not text.strip(),
            "fallback_reason": (
                "PDF_NATIVE_TEXT_EMPTY" if suffix == ".pdf" and not text.strip() else None
            ),
            "trust_boundary": "untrusted_external_content",
            "prompt_injection_signals": self._prompt_injection_signals(text),
            "macros_executed": False,
        }

    def vision_inputs(
        self,
        value: str,
        *,
        pages: list[int] | None = None,
        max_pages: int = 10,
    ) -> list[dict[str, Any]]:
        path = self._safe_file(value)
        if self._is_image(path):
            if path.stat().st_size > 20 * 1024 * 1024:
                raise ValueError("Vision image exceeds the 20 MiB input bound")
            return [
                {
                    "label": path.name,
                    "media_type": mimetypes.guess_type(path.name)[0] or "image/png",
                    "content": path.read_bytes(),
                }
            ]
        if path.suffix.casefold() != ".pdf":
            raise ValueError("Vision fallback accepts an image or PDF")
        import fitz

        document = fitz.open(path)
        try:
            indexes = self._page_indexes(document.page_count, pages)[:max_pages]
            return [
                {
                    "label": f"{path.name} page {index + 1}",
                    "media_type": "image/png",
                    "content": document[index]
                    .get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
                    .tobytes("png"),
                }
                for index in indexes
            ]
        finally:
            document.close()

    def list_archive(self, value: str, *, limit: int = 1_000) -> dict[str, Any]:
        path = self._safe_file(value)
        if not zipfile.is_zipfile(path):
            raise ValueError("Archive listing currently supports ZIP/Office containers")
        entries: list[dict[str, Any]] = []
        safe = True
        total = 0
        with zipfile.ZipFile(path) as archive:
            for info in archive.infolist()[:limit]:
                name_path = Path(info.filename)
                entry_safe = not (
                    name_path.is_absolute()
                    or ".." in name_path.parts
                    or info.file_size > 512 * 1024 * 1024
                    or (info.compress_size and info.file_size / info.compress_size > 1_000)
                )
                safe = safe and entry_safe
                total += info.file_size
                entries.append(
                    {
                        "name": info.filename,
                        "size": info.file_size,
                        "compressed_size": info.compress_size,
                        "safe": entry_safe,
                        "directory": info.is_dir(),
                    }
                )
        return {
            **self.metadata(path),
            "entries": entries,
            "entry_count": len(entries),
            "uncompressed_size": total,
            "safe": safe,
            "executed": False,
        }

    def copy(self, source: str, destination: str) -> dict[str, object]:
        source_path = self._resolve(source, must_exist=True)
        destination_path = self._resolve(destination, must_exist=False)
        self._deny_secret(source_path)
        self._prepare_destination(destination_path)
        shutil.copy2(source_path, destination_path)
        return self.metadata(destination_path)

    def move(self, source: str, destination: str) -> dict[str, object]:
        source_path = self._resolve(source, must_exist=True)
        destination_path = self._resolve(destination, must_exist=False)
        self._deny_secret(source_path)
        self._prepare_destination(destination_path)
        shutil.move(str(source_path), str(destination_path))
        return self.metadata(destination_path)

    def rename(self, source: str, name: str) -> dict[str, object]:
        if Path(name).name != name or name in {".", ".."}:
            raise ValueError("Rename accepts a filename, not a path")
        source_path = self._resolve(source, must_exist=True)
        return self.move(str(source_path), str(source_path.with_name(name)))

    def delete_to_trash(self, value: str) -> dict[str, object]:
        path = self._resolve(value, must_exist=True)
        self._deny_secret(path)
        self.trash_root.mkdir(parents=True, exist_ok=True, mode=0o700)
        target = self.trash_root / f"{uuid4()}-{path.name}"
        shutil.move(str(path), str(target))
        return {
            "original_path": str(path),
            "trash_path": str(target),
            "recoverable": True,
        }

    def metadata(self, path: Path) -> dict[str, object]:
        resolved = path.resolve(strict=True)
        self._assert_allowed(resolved)
        stat = resolved.stat()
        return {
            "path": str(resolved),
            "name": resolved.name,
            "size": stat.st_size,
            "modified_at": stat.st_mtime,
            "media_type": mimetypes.guess_type(resolved.name)[0],
        }

    def _safe_file(self, value: str) -> Path:
        path = self._resolve(value, must_exist=True)
        self._deny_secret(path)
        if not path.is_file():
            raise ValueError("Path is not a regular file")
        if path.suffix.casefold() in {".docm", ".xlsm", ".pptm"}:
            raise PermissionError("Macro-enabled Office files are not parsed")
        return path

    @staticmethod
    def _is_image(path: Path) -> bool:
        media_type = mimetypes.guess_type(path.name)[0] or ""
        return media_type.startswith("image/")

    @staticmethod
    def _page_indexes(count: int, pages: list[int] | None) -> list[int]:
        if not pages:
            return list(range(count))
        indexes = sorted(set(page - 1 for page in pages))
        if any(index < 0 or index >= count for index in indexes):
            raise ValueError("Requested PDF page is outside the document")
        return indexes

    @staticmethod
    def _prompt_injection_signals(text: str) -> list[str]:
        normalized = text.casefold()
        signals = []
        patterns = {
            "IGNORE_PREVIOUS_INSTRUCTIONS": "ignore previous instructions",
            "SYSTEM_PROMPT_REQUEST": "system prompt",
            "TOOL_PERMISSION_REQUEST": "grant permission",
            "JAPANESE_OVERRIDE_REQUEST": "以前の指示を無視",
        }
        for code, pattern in patterns.items():
            if pattern in normalized:
                signals.append(code)
        return signals

    def _resolve(self, value: str, *, must_exist: bool) -> Path:
        candidate = Path(value).expanduser()
        if not candidate.is_absolute():
            if len(self.roots) != 1:
                raise ValueError("Relative paths require exactly one configured file root")
            candidate = self.roots[0] / candidate
        resolved = candidate.resolve(strict=False)
        self._assert_allowed(resolved)
        if must_exist and not resolved.exists():
            raise FileNotFoundError(resolved)
        return resolved

    def _assert_allowed(self, path: Path) -> None:
        if not self.roots or not any(
            path == root or path.is_relative_to(root) for root in self.roots
        ):
            raise PermissionError("Path is outside configured file roots")

    @staticmethod
    def _prepare_destination(path: Path) -> None:
        if path.exists():
            raise FileExistsError("Destination already exists")
        path.parent.mkdir(parents=True, exist_ok=True)

    @staticmethod
    def _is_secret(path: Path) -> bool:
        lowered_parts = {part.casefold() for part in path.parts}
        return (
            path.name.casefold() in _SECRET_NAMES
            or path.suffix.casefold() in _SECRET_SUFFIXES
            or ".ssh" in lowered_parts
            or ".gnupg" in lowered_parts
        )

    @classmethod
    def _deny_secret(cls, path: Path) -> None:
        if cls._is_secret(path):
            raise PermissionError("Secret/key files are excluded from agent file access")
